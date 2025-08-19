from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from lxml import etree
import json
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from math import cos, sin, pi
from Set_Text import TextRunFactory
from utils import textwrap, requests
from config import ollama_url, used_model

def add_l_connector(slide, start, bend, end): 
    """ æ’å…¥ L å‹é€£ç·šï¼ˆå…©æ®µï¼‰ä¸¦çµ„æˆç¾¤çµ„ï¼Œå°¾ç«¯åŠ ä¸Šç®­é ­ """
    left1, top1 = start[0], start[1]
    left2, top2 = bend[0], bend[1]
    left3, top3 = end[0], end[1]

    # æ’å…¥å…©æ®µ connectorï¼šstart â†’ bend â†’ end
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left1, top1, left2, top2)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left3, top3, left2, top2)

    for line in [line1, line2]:
        line.line.width = Pt(2)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.end_arrowhead = None  # å…ˆæ¸…æ‰é è¨­ç®­é ­

    # åªåœ¨å°¾ç«¯ line2 åŠ ç®­é ­ï¼ˆXML æ“ä½œï¼‰
    line_elem = line2.line._get_or_add_ln()
    line_elem.append(parse_xml(
        '<a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    ))

    # å°‡å…©æ®µç·šçµ„æˆä¸€å€‹ç¾¤çµ„
    group = slide.shapes.add_group_shape([line1, line2])
    return group

def draw_connectors(slide, nodes, shapes, layout_key):
    for node_id, node in nodes.items():
        start_shape = shapes[node_id]
        for nxt_id in node["next"]:
            if nxt_id not in shapes:
                continue
            end_shape = shapes[nxt_id]

            # åŸºæœ¬å¹¾ä½•å®šä½é»
            start_top_y = start_shape.top
            start_center_x = start_shape.left + start_shape.width / 2
            start_center_y = start_shape.top + start_shape.height / 2
            start_bottom_y = start_shape.top + start_shape.height
            start_right_x = start_shape.left + start_shape.width
            start_left_x = start_shape.left
            
            end_top_y = end_shape.top
            end_center_x = end_shape.left + end_shape.width / 2
            end_center_y = end_shape.top + end_shape.height / 2
            end_left_x = end_shape.left
            end_right_x = end_shape.left + end_shape.width

            dx = end_shape.left - start_shape.left
            dy = end_shape.top - start_shape.top

            arrow_margin = int(start_shape.width/5)

            # æ¢ä»¶ 1ï¼šA åœ¨ B æ­£ä¸Šæ–¹ï¼ˆx åº§æ¨™å·®å°ï¼Œy åº§æ¨™å·®ç‚ºæ­£ï¼‰
            if abs(dx) < 20 and dy > 0:
                # â†“ å¾ A åº•éƒ¨ â†’ B é ‚éƒ¨
                start_x = int(start_center_x)
                start_y = int(start_bottom_y)
                end_x = int(end_center_x)
                end_y = int(end_top_y)
                
            # æ¢ä»¶ 2ï¼šA åœ¨ B æ­£å·¦æ–¹ï¼ˆy å·®å°ï¼Œx å·®ç‚ºæ­£ï¼‰
            elif abs(dy) < 20 and dx > 0:
                # â†’ å¾ A å³å´ â†’ B å·¦å´
                start_x = int(start_right_x)
                start_y = int(start_center_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
			
			# æ¢ä»¶ 4ï¼šA åœ¨ B å³å´ï¼ˆæ–°å¢ï¼‰
            elif abs(dy) < 20 and dx < 0:
                # â† å¾ A å·¦å´ â†’ B å³å´
                start_x = int(start_left_x)
                start_y = int(start_center_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)
                #print(node_id, nxt_id, "å·¦")

            # æ¢ä»¶ 3ï¼šA åœ¨ B å·¦ä¸Šæ–¹ï¼Œä¸” B å·¦é‚Š > A ä¸­å¿ƒï¼ˆä»£è¡¨ B ç¨å¾€å³åï¼‰
            elif dx > 0 and dy > 0 and end_left_x > start_center_x:
                # â†˜ å¾ A åº•éƒ¨ â†’ B å·¦å´
                start_x = int(start_center_x)+arrow_margin
                start_y = int(start_bottom_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "å³ä¸‹")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            # æ¢ä»¶ 6ï¼šA åœ¨ B å·¦ä¸‹æ–¹
            elif dx > 0 and dy < 0 :
                # å¾ A ä¸Šæ–¹ â†’ B å·¦å´
                start_x = int(start_center_x)+arrow_margin
                start_y = int(start_top_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "å³ä¸Š")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
			# æ¢ä»¶ 5ï¼šA åœ¨ B å³ä¸Šæ–¹
            elif dx < 0 and dy > 0 :
                # â†™ å¾ A åº•éƒ¨ â†’ B å³å´
                start_x = int(start_center_x)-arrow_margin
                start_y = int(start_bottom_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)
                #print(node_id, nxt_id, "å·¦ä¸‹")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            # æ¢ä»¶ 7ï¼šA åœ¨ B å³ä¸‹æ–¹
            elif dx < 0 and dy < 0 :
                # â†™ å¾ A åº•éƒ¨ â†’ B å³å´
                start_x = int(start_center_x)-arrow_margin
                start_y = int(start_top_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)
                #print(node_id, nxt_id, "å·¦ä¸Š")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            else:
                # é è¨­ fallbackï¼šå¾ A å³å´ â†’ B å·¦å´
                start_x = int(start_right_x)
                start_y = int(start_center_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "ä¸€èˆ¬")

            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.ELBOW, end_x, end_y, start_x, start_y
            )
            connector.line.color.rgb = RGBColor(0, 0, 0)
            connector.line.width = Pt(2)
            line_elem = connector.line._get_or_add_ln()
            line_elem.append(parse_xml(
                '<a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))

# æª¢æŸ¥å…©å€‹æ¡†æ¡†æ˜¯å¦é‡ç–Šï¼Œä¿ç•™ margin é–“è·
def is_overlap(rect1, rect2, margin=10):
    x1, y1, w1, h1 = rect1
    x2, y2, w2, h2 = rect2

    return not (
        x1 + w1 + margin <= x2 or  # rect1 åœ¨ rect2 å·¦é‚Š + margin
        x2 + w2 + margin <= x1 or  # rect2 åœ¨ rect1 å·¦é‚Š + margin
        y1 + h1 + margin <= y2 or  # rect1 åœ¨ rect2 ä¸Šæ–¹ + margin
        y2 + h2 + margin <= y1     # rect2 åœ¨ rect1 ä¸Šæ–¹ + margin
    )

# å»ºç«‹æŠ•å½±ç‰‡èˆ‡é¿å…é‡ç–Šçš„è‡ªå‹•æ’ç‰ˆ
def create_slide(prs, nodes, layout_key, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}
    shape_positions = []  # å„²å­˜æ‰€æœ‰å·²æ“ºæ”¾çš„æ¡†æ¡†ä½ç½® (x, y, w, h)
    H_START = 100

    for node_id, node in nodes.items():
        layout = node["layouts"][layout_key]
        x = Inches(layout["x"] * 1.4 / 100)
        y = Inches((H_START + layout["y"] * 1.6) / 100)
        w = Inches(layout["width"] * 1.2 / 100)
        h = Inches(layout["height"] * 1.6 / 100)

        # å˜—è©¦é¿å…èˆ‡ç¾æœ‰æ¡†æ¡†é‡ç–Š
        max_attempts = 100
        dy = Inches(0.15)  # æ¯æ¬¡å¾®èª¿çš„è·é›¢
        attempt = 0
        while any(is_overlap((x, y, w, h), rect) for rect in shape_positions) and attempt < max_attempts:
            y += dy
            attempt += 1

        # åŠ å…¥æ¡†æ¡†
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(253, 234, 218)
        shape.line.color.rgb = RGBColor(0, 0, 0)

        # æ–‡å­—æ¨£å¼
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.clear()

        run_factory = TextRunFactory(p)
        run_factory.add_icon(node.get("icon", ""), font_size=35, color=RGBColor(152, 72, 7))
        run_factory.add_id(node["id"])

        shapes[node_id] = shape
        shape_positions.append((x, y, w, h))  # è¨˜éŒ„æ­¤æ¡†æ¡†ä½ç½®

    draw_connectors(slide, nodes, shapes, layout_key)

def create_list_slide(prs, nodes, title_text="æ¸…å–®åœ–"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}

    num_nodes = len(nodes)
    spacing = Inches(0.2)
    total_width = prs.slide_width - Inches(0.6)
    shape_width = (total_width - spacing * (num_nodes - 1)) / num_nodes
    shape_height = Inches(1)
    top_margin = Inches(2)

    for i, (node_id, node) in enumerate(nodes.items()):
        left = Inches(0.3) + i * (shape_width + spacing)

        shape_1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin, shape_width, shape_height)
        shape_1.fill.background()  # è®“å¡«è‰²é€æ˜
        shape_1.line.fill.background()  # è®“é‚Šæ¡†ä¹Ÿé€æ˜ï¼ˆå®Œå…¨ä¸é¡¯ç¤ºï¼‰
        tf_1 = shape_1.text_frame
        tf_1.clear()

        # ç¬¬ä¸€è¡Œ icon
        p_1 = tf_1.paragraphs[0]
        p_1.alignment = PP_ALIGN.CENTER
        p_1.clear()

        run_factory = TextRunFactory(p_1)
        run_factory.add_icon(node.get("icon", ""))

        shapes[node_id] = shape_1

        text_len = len(node["id"]+node["add"])
        if text_len < 50:
            shape_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin + shape_height, shape_width, shape_height*2)
        else:
            shape_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin + shape_height, shape_width, shape_height*3)
        # è¨­å®šå¡«è‰²èˆ‡é‚Šæ¡†
        shape_2.fill.solid()
        shape_2.fill.fore_color.rgb = RGBColor(218, 235, 247)
        shape_2.line.color.rgb = RGBColor(0, 0, 0)

        tf_2 = shape_2.text_frame
        tf_2.clear()

        p_2 = tf_2.paragraphs[0]
        p_2.alignment = PP_ALIGN.CENTER
        p_2.space_after = Pt(4)
        p_2.clear()

        # ç¬¬äºŒè¡Œ ID, ç¬¬ä¸‰è¡Œ Add
        run_factory_2 = TextRunFactory(p_2)
        run_factory_2.add_id(node["id"])
        run_factory_2.add_add(node["add"])

def create_cycle_slide(prs, nodes, title_text="å¾ªç’°åœ–"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}

    # åœ“çš„åƒæ•¸
    num_nodes = len(nodes)
    radius = Inches(1.5 + num_nodes * 0.1)
    center_x = prs.slide_width / 2
    center_y = prs.slide_height / 2 + Inches(0.5)

    shape_w = Inches(2)
    shape_h = Inches(1.8)

    for i, (node_id, node) in enumerate(nodes.items()):
        angle = 2 * pi * i / num_nodes - pi / 2  # æ¯å€‹ç¯€é»çš„è§’åº¦
        x = center_x + radius * cos(angle) - shape_w / 2
        y = center_y + radius * sin(angle) - shape_h / 2

        # åœ“ä¸Šçš„æ¡†æ¡†
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, shape_w, shape_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(221, 235, 247)
        shape.line.color.rgb = RGBColor(0, 0, 0)

        tf = shape.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.clear()

        # ç¬¬ä¸€è¡Œ icon, ç¬¬äºŒè¡Œ ID, ç¬¬ä¸‰è¡Œ Add
        run_factory = TextRunFactory(p)
        run_factory.add_icon(node.get("icon", ""))
        run_factory.add_id(node["id"])

        

def detect_layout_types(nodes):
    for node in nodes.values():
        return list(node.get("layouts", {}).keys())[:2]
    return []

def create_node(summary_text):
    prompt = textwrap.dedent(f"""
    è«‹æ ¹æ“šä»¥ä¸‹çš„æ ¼å¼ï¼Œèƒå–å‡ºå…¶ä¸­çš„ä¸»è¦æµç¨‹æˆ–æ¶æ§‹ï¼Œä¸¦å°‡å…¶æ‹†è§£ç‚ºçµæ§‹åŒ–çš„ç¯€é»ä¿¡æ¯ã€‚æ¯å€‹ç¯€é»æ‡‰è©²åŒ…å«ä»¥ä¸‹æ¬„ä½ï¼š
    - **id**: ç¯€é»çš„å”¯ä¸€è­˜åˆ¥å­—ä¸²ï¼Œé€™å°‡æ˜¯è©²ç¯€é»çš„é¡¯ç¤ºæ–‡å­—ï¼ˆä¸é‡è¤‡ï¼Œé©ç•¶æ™‚å¯ä»¥ä½¿ç”¨ä¸­æ–‡ï¼‰ã€‚
    - **add**: æ­¤ç¯€é»çš„ç°¡è¦èªªæ˜, æœ€å¤šä¸è¶…é20å­—ã€‚
    - **next**: é€™æ˜¯ä¸€å€‹åˆ—è¡¨ï¼ŒåŒ…å«ç•¶å‰ç¯€é»çš„æŒ‡å‘ç¯€é»ï¼ˆå³å¾ŒçºŒæ­¥é©Ÿï¼‰ã€‚å¦‚æœæ­¤ç¯€é»æ˜¯æµç¨‹çš„æœ€çµ‚æ­¥é©Ÿï¼Œå‰‡è¨­ç‚ºç©ºåˆ—è¡¨ []ã€‚
    - **icon**: é€™æ˜¯èˆ‡ç¯€é»ç›¸é—œçš„ Unicode è¡¨æƒ…ç¬¦è™Ÿï¼Œç”¨ä¾†è¡¨ç¤ºè©²ç¯€é»çš„è±¡å¾µæ„åœ–æˆ–åŠŸèƒ½ã€‚

    è«‹æ³¨æ„ï¼š
    - è‹¥ç„¡æ˜ç¢ºçµæ§‹æˆ–å±¤ç´šï¼Œè«‹åˆç†æ¨æ¸¬åˆ†å±¤ä¸¦æ‰¾å‡ºé—œè¯æ€§ã€‚
    - æ ¹æ“šæµç¨‹æè¿°ï¼Œæ‹†è§£æˆ 2 åˆ° 4 å€‹ä¸»è¦ç¯€é»ï¼Œä¸¦åœ¨é€™äº›ç¯€é»ä¹‹é–“å»ºç«‹ã€Œä¸‹ä¸€æ­¥ã€çš„é—œè¯ã€‚
    - ä¿è­‰æ¯å€‹ç¯€é»éƒ½åŒ…å«ä»¥ä¸Šæ‰€åˆ—çš„æ¬„ä½ã€‚

    ### ç¯„ä¾‹ï¼š
    {{
        "ç¯€é»1": {{
            "id": "è‡ªå‹•åŒ–æ ¸å¿ƒæ¦‚å¿µ",
            "add": "",
            "next": ["è§¸ç™¼ç¯€é»", "å‹•ä½œç¯€é»"],
            "icon": "ğŸ”"
        }},
        "ç¯€é»2": {{
            "id": "è§¸ç™¼ç¯€é»",
            "add": "ä½•æ™‚åŸ·è¡Œå·¥ä½œæµç¨‹ï¼Ÿ",
            "next": ["è³‡æ–™å„²å­˜"],
            "icon": "ğŸ§©"
        }}
    }}

    ### ä»¥ä¸‹æ˜¯è¦è™•ç†çš„æ–‡å­—ï¼š
    {summary_text}
    """)
    
    response = requests.post(
        f"{ollama_url}/api/generate",
        json={"model": used_model, "prompt": prompt, "stream": False}
    )

    # print(f"åŸå§‹å»ºé€ nodeçš„å›æ‡‰: {response.text}")     # debugging line

    if response.status_code == 200:
        # If the request is successful, parse the response
        node_data = response.json().get("response", "")
        
        # Remove the ```json\n``` and closing ``` part
        if node_data:
            clean_node_data = node_data.replace("```json\n", "").replace("```", "").strip()
            
            # Now try to parse the clean JSON data
            try:
                nodes = json.loads(clean_node_data)
                return nodes  # Return the extracted nodes as JSON
            except json.JSONDecodeError:
                print("Error: Failed to decode JSON from cleaned response.")
                return None
        else:
            print("Error: No response data found.")
            return None
    else:
        print(f"Error: API request failed with status code {response.status_code}")
        return None

def generate_diagram_to_ppt(save_path, st_status, node_data):
    print(node_data)
    msg = "ğŸ“Š è£½ä½œåœ–è¡¨ä¸­..."
    if st_status:
        st_status.info(msg)

    prs = Presentation(save_path)

    nodes = {n["id"]: n for n in node_data.values()}  # æ­£ç¢ºè½‰æ› key

    layout_types = detect_layout_types(nodes)
    for layout in layout_types:
        create_slide(prs, nodes, layout, layout)
    create_list_slide(prs, nodes, "æ¸…å–®åœ–")
    create_cycle_slide(prs, nodes, "å¾ªç’°åœ–")

    msg = "ğŸ“Š è£½ä½œå®Œæˆ..."
    if st_status:
        st_status.info(msg)
    prs.save(save_path)      # å„²å­˜ç‚º PPT æª”æ¡ˆ
